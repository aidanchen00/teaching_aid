#!/usr/bin/env python3
"""
PowerPoint Generation Script using python-pptx
Generates business strategy presentations from JSON input
"""

import json
import sys
import base64
from io import BytesIO
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    print(json.dumps({"error": f"python-pptx not installed or import error: {str(e)}. Run: pip install python-pptx"}))
    sys.exit(1)


def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def create_presentation(data: dict) -> bytes:
    """Create a PowerPoint presentation from structured data"""

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Extract branding
    branding = data.get("branding", {})
    primary_color = branding.get("primaryColor", "#3B82F6")
    secondary_color = branding.get("secondaryColor", "#10B981")
    accent_color = branding.get("accentColor", "#F59E0B")

    primary_rgb = hex_to_rgb(primary_color)
    secondary_rgb = hex_to_rgb(secondary_color)
    accent_rgb = hex_to_rgb(accent_color)

    company_name = data.get("companyName", "Company")

    slides_data = data.get("slides", [])

    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            create_title_slide(prs, slide_data, primary_rgb, secondary_rgb, company_name)
        elif slide_type == "section":
            create_section_slide(prs, slide_data, primary_rgb, accent_rgb)
        elif slide_type == "content":
            create_content_slide(prs, slide_data, primary_rgb, secondary_rgb)
        elif slide_type == "two_column":
            create_two_column_slide(prs, slide_data, primary_rgb, secondary_rgb)
        elif slide_type == "metrics":
            create_metrics_slide(prs, slide_data, primary_rgb, secondary_rgb, accent_rgb)
        elif slide_type == "timeline":
            create_timeline_slide(prs, slide_data, primary_rgb, secondary_rgb, accent_rgb)
        elif slide_type == "closing":
            create_closing_slide(prs, slide_data, primary_rgb, secondary_rgb, company_name)
        else:
            create_content_slide(prs, slide_data, primary_rgb, secondary_rgb)

    # Save to bytes
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)

    return pptx_buffer.getvalue()


def add_branded_header(slide, title: str, primary_rgb: tuple):
    """Add a branded header bar to the slide"""
    # Header background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*primary_rgb)
    shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def create_title_slide(prs, slide_data: dict, primary_rgb: tuple, secondary_rgb: tuple, company_name: str):
    """Create a title slide"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Background gradient effect (solid for simplicity)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(*primary_rgb)
    background.line.fill.background()

    # Accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.5),
        Inches(13.333), Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(*secondary_rgb)
    accent_bar.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5),
        Inches(11.5), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "Business Strategy")
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.2),
        Inches(11.5), Inches(0.8)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("subtitle", company_name)
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Date/tagline
    if slide_data.get("tagline"):
        tagline_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.2),
            Inches(11.5), Inches(0.5)
        )
        tf = tagline_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("tagline", "")
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.alignment = PP_ALIGN.CENTER


def create_section_slide(prs, slide_data: dict, primary_rgb: tuple, accent_rgb: tuple):
    """Create a section divider slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Left colored bar
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.3), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = RGBColor(*primary_rgb)
    left_bar.line.fill.background()

    # Section number/icon area
    icon_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1), Inches(2.8),
        Inches(1.5), Inches(1.5)
    )
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = RGBColor(*accent_rgb)
    icon_circle.line.fill.background()

    # Section number
    if slide_data.get("sectionNumber"):
        num_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.1),
            Inches(1.5), Inches(1)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_data.get("sectionNumber", ""))
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(3), Inches(2.8),
        Inches(9), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "Section")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*primary_rgb)

    # Section description
    if slide_data.get("description"):
        desc_box = slide.shapes.add_textbox(
            Inches(3), Inches(4.2),
            Inches(9), Inches(1)
        )
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("description", "")
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(100, 100, 100)


def create_content_slide(prs, slide_data: dict, primary_rgb: tuple, secondary_rgb: tuple):
    """Create a standard content slide with bullets"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_branded_header(slide, slide_data.get("title", ""), primary_rgb)

    # Content area
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(12.3), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    bullets = slide_data.get("bullets", [])
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(12)

    # Optional subtitle/key point
    if slide_data.get("keyPoint"):
        key_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.3),
            Inches(12.3), Inches(0.8)
        )
        key_box.fill.solid()
        key_box.fill.fore_color.rgb = RGBColor(*secondary_rgb)
        key_box.fill.fore_color.brightness = 0.8
        key_box.line.fill.background()

        key_text = slide.shapes.add_textbox(
            Inches(0.7), Inches(6.45),
            Inches(12), Inches(0.5)
        )
        tf = key_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"Key Insight: {slide_data.get('keyPoint')}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*secondary_rgb)


def create_two_column_slide(prs, slide_data: dict, primary_rgb: tuple, secondary_rgb: tuple):
    """Create a two-column layout slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_branded_header(slide, slide_data.get("title", ""), primary_rgb)

    # Left column
    left_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(5.8), Inches(0.5)
    )
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("leftTitle", "")
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*primary_rgb)

    left_content = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.1),
        Inches(5.8), Inches(4.5)
    )
    tf = left_content.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(slide_data.get("leftBullets", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(8)

    # Divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.5), Inches(1.5),
        Inches(0.02), Inches(5.5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(200, 200, 200)
    divider.line.fill.background()

    # Right column
    right_title = slide.shapes.add_textbox(
        Inches(7), Inches(1.5),
        Inches(5.8), Inches(0.5)
    )
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("rightTitle", "")
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*secondary_rgb)

    right_content = slide.shapes.add_textbox(
        Inches(7), Inches(2.1),
        Inches(5.8), Inches(4.5)
    )
    tf = right_content.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(slide_data.get("rightBullets", [])):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(8)


def create_metrics_slide(prs, slide_data: dict, primary_rgb: tuple, secondary_rgb: tuple, accent_rgb: tuple):
    """Create a metrics/KPI slide with boxes"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_branded_header(slide, slide_data.get("title", "Key Metrics"), primary_rgb)

    metrics = slide_data.get("metrics", [])
    colors = [primary_rgb, secondary_rgb, accent_rgb, (139, 92, 246)]  # Add purple

    num_metrics = min(len(metrics), 4)
    box_width = 2.8
    total_width = num_metrics * box_width + (num_metrics - 1) * 0.3
    start_x = (13.333 - total_width) / 2

    for i, metric in enumerate(metrics[:4]):
        x = start_x + i * (box_width + 0.3)
        color = colors[i % len(colors)]

        # Metric box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.2),
            Inches(box_width), Inches(3.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*color)
        box.line.fill.background()

        # Metric value
        value_box = slide.shapes.add_textbox(
            Inches(x), Inches(2.8),
            Inches(box_width), Inches(1.2)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(metric.get("value", ""))
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Metric label
        label_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(4.1),
            Inches(box_width - 0.2), Inches(0.8)
        )
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = metric.get("label", "")
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Metric trend/description
        if metric.get("trend"):
            trend_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(4.9),
                Inches(box_width - 0.2), Inches(0.5)
            )
            tf = trend_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric.get("trend", "")
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(220, 220, 220)
            p.alignment = PP_ALIGN.CENTER


def create_timeline_slide(prs, slide_data: dict, primary_rgb: tuple, secondary_rgb: tuple, accent_rgb: tuple):
    """Create a timeline/roadmap slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_branded_header(slide, slide_data.get("title", "Roadmap"), primary_rgb)

    milestones = slide_data.get("milestones", [])
    num_milestones = min(len(milestones), 5)

    if num_milestones == 0:
        return

    # Timeline line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(4),
        Inches(11.333), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*primary_rgb)
    line.line.fill.background()

    spacing = 11.333 / num_milestones
    colors = [primary_rgb, secondary_rgb, accent_rgb, (139, 92, 246), (236, 72, 153)]

    for i, milestone in enumerate(milestones[:5]):
        x = 1 + (i * spacing) + (spacing / 2) - 0.4
        color = colors[i % len(colors)]

        # Circle marker
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(3.75),
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*color)
        circle.line.fill.background()

        # Time label (above)
        time_box = slide.shapes.add_textbox(
            Inches(x - 0.5), Inches(2.8),
            Inches(1.5), Inches(0.6)
        )
        tf = time_box.text_frame
        p = tf.paragraphs[0]
        p.text = milestone.get("time", "")
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*color)
        p.alignment = PP_ALIGN.CENTER

        # Milestone text (below)
        text_box = slide.shapes.add_textbox(
            Inches(x - 0.8), Inches(4.5),
            Inches(2), Inches(1.5)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = milestone.get("title", "")
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.alignment = PP_ALIGN.CENTER

        if milestone.get("description"):
            p = tf.add_paragraph()
            p.text = milestone.get("description", "")
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(120, 120, 120)
            p.alignment = PP_ALIGN.CENTER


def create_closing_slide(prs, slide_data: dict, primary_rgb: tuple, secondary_rgb: tuple, company_name: str):
    """Create a closing/thank you slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(*primary_rgb)
    background.line.fill.background()

    # Main message
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5),
        Inches(11.5), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "Thank You")
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Contact/CTA
    if slide_data.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(4.2),
            Inches(11.5), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("subtitle", "")
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.alignment = PP_ALIGN.CENTER

    # Contact details
    if slide_data.get("contact"):
        contact_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(5.5),
            Inches(11.5), Inches(1)
        )
        tf = contact_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("contact", "")
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.CENTER


def main():
    """Main entry point - reads JSON from stdin, outputs base64 PPTX"""
    try:
        # Read JSON input from stdin
        input_data = sys.stdin.read()
        data = json.loads(input_data)

        # Generate presentation
        pptx_bytes = create_presentation(data)

        # Output base64 encoded result
        result = {
            "success": True,
            "pptxBase64": base64.b64encode(pptx_bytes).decode('utf-8'),
            "slideCount": len(data.get("slides", []))
        }

        print(json.dumps(result))

    except json.JSONDecodeError as e:
        print(json.dumps({"success": False, "error": f"Invalid JSON: {str(e)}"}))
        sys.exit(1)
    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}))
        sys.exit(1)


if __name__ == "__main__":
    main()
